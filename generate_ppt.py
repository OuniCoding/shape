#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
千里風光計畫 — 完整 PPT 生成腳本
包含：6 張圖表 + 12 頁幻燈片
"""

import os
import warnings
import matplotlib
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ==================== 全局設定 ====================
warnings.filterwarnings('ignore', message='Glyph.*missing from current font')
matplotlib.use('Agg')

# 自動尋找中文字體
def find_chinese_font():
    candidates = [
        'SimHei', 'Microsoft YaHei', 'WenQuanYi Micro Hei',
        'Noto Sans CJK SC', 'Noto Sans CJK TC', 'Noto Sans Mono CJK SC',
        'PingFang SC', 'PingFang TC', 'Heiti SC', 'Heiti TC',
        'STHeiti', 'STSong', 'Source Han Sans SC', 'Source Han Sans TC',
        'Arial Unicode MS'
    ]
    available_fonts = set(f.name for f in fm.fontManager.ttflist)
    for font in candidates:
        if font in available_fonts:
            print(f"[INFO] 使用中文字體: {font}")
            return font
    print("[WARN] 未找到中文字體，圖表中文可能顯示為方框")
    return 'DejaVu Sans'

chinese_font = find_chinese_font()
plt.rcParams['font.sans-serif'] = [chinese_font, 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

# 輸出路徑
output_dir = './千里風光計畫_PPT/'
os.makedirs(output_dir, exist_ok=True)

# 顏色定義
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_BLUE = RGBColor(0x2B, 0x6C, 0xB4)
ACCENT_GREEN = RGBColor(0x2E, 0x8B, 0x57)
ACCENT_ORANGE = RGBColor(0xD4, 0x69, 0x2A)
ACCENT_RED = RGBColor(0xC4, 0x1E, 0x3A)
ACCENT_PURPLE = RGBColor(0x6B, 0x4C, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ==================== 圖表 1：組織架構圖 ====================
fig, ax = plt.subplots(figsize=(10, 5))
ax.set_xlim(0, 10)
ax.set_ylim(0, 5)
ax.axis('off')

# 母公司
rect1 = plt.Rectangle((3.5, 3.8), 3, 0.8, facecolor='#2B6CB4', edgecolor='#1A1A1A', linewidth=2, alpha=0.95)
ax.add_patch(rect1)
# ax.text(5, 4.2, 'Group Parent\nRegional Smart Grid Asset Holding', ha='center', va='center',
#         fontsize=11, color='white', weight='bold')
ax.text(5, 4.2, '集團母公司\n區域級智慧電網資產控股與營運平台', ha='center', va='center', fontsize=11, color='white', weight='bold')
# 連接線
ax.plot([5, 5], [3.8, 3.2], color='#666666', linewidth=2)
ax.plot([1.5, 8.5], [3.2, 3.2], color='#666666', linewidth=2)
ax.plot([1.5, 1.5], [3.2, 2.8], color='#666666', linewidth=2)
ax.plot([5, 5], [3.2, 2.8], color='#666666', linewidth=2)
ax.plot([8.5, 8.5], [3.2, 2.8], color='#666666', linewidth=2)

# 子公司
# subs = [
#     (1.5, 'Sub 1: EPC', 'Regional Engineering Delivery', '#E3F2FD', '#2B6CB4'),
#     (5, 'Sub 2: Investment/Dev', 'Group Core - Asset Operator', '#E8F5E9', '#2E8B57'),
#     (8.5, 'Sub 3: Trading', 'Energy Market Expansion Engine', '#FFF3E0', '#D4692A'),
# ]
subs = [
    (1.5, 'EPC 工程總包公司', '區域工程交付平台', '#E3F2FD', '#2B6CB4'),
    (5, '投資/開發/方案公司', '集團真正核心', '#E8F5E9', '#2E8B57'),
    (8.5, '貿易公司', '能源市場擴張引擎', '#FFF3E0', '#D4692A'),
]
for x, t, d, fc, ec in subs:
    rect = plt.Rectangle((x - 1.3, 1.8), 2.6, 0.9, facecolor=fc, edgecolor=ec, linewidth=2)
    ax.add_patch(rect)
    ax.text(x, 2.35, t, ha='center', va='center', fontsize=10, color=ec, weight='bold')
    ax.text(x, 2.0, d, ha='center', va='center', fontsize=8, color='#666666')

# 底部標註
# ax.text(1.5, 1.6, 'Rapid Deployment\nAsset Control', ha='center', va='top', fontsize=8, color='#666666')
# ax.text(5, 1.6, 'Hold Assets\nVPP & Cash Flow', ha='center', va='top', fontsize=8, color='#666666')
# ax.text(8.5, 1.6, 'Market Expansion\nCustomer Acquisition', ha='center', va='top', fontsize=8, color='#666666')
ax.text(1.5, 1.6, '快速部署與資產控制', ha='center', va='top', fontsize=8, color='#666666')
ax.text(5, 1.6, '持有資產、VPP和現金流', ha='center', va='top', fontsize=8, color='#666666')
ax.text(8.5, 1.6, '市場拓展與獲取客戶', ha='center', va='top', fontsize=8, color='#666666')

plt.tight_layout()
plt.savefig(f'{output_dir}chart_org.png', dpi=150, bbox_inches='tight', facecolor='white')
plt.close()
print("[OK] chart_org.png")

# ==================== 圖表 2：EPC 收入結構 ====================
fig, ax = plt.subplots(figsize=(6, 4))
# labels = ['EPC 40%', 'O&M 30%', 'Upgrade 20%', 'AI 10%']
labels = ['EPC 施工\n40%', '運維\n30%', '改造升級\n20%', 'AI 部署\n10%']
sizes = [40, 30, 20, 10]
colors = ['#2B6CB4', '#2E8B57', '#D4692A', '#6B4C9A']
# wedges, texts = ax.pie(sizes, colors=colors, startangle=90,
#                         textprops={'fontsize': 10, 'color': 'white', 'weight': 'bold'})
# for i, text in enumerate(texts):
#     text.set_text(labels[i])
wedges, texts, autotexts = ax.pie(
    sizes,
    # explode=explode,
    labels=labels,              # ← 必須在這裡傳入 labels
    colors=colors,
    startangle=90,
    autopct='%1.0f%%',          # ← 自動顯示百分比
    pctdistance=0.6,            # ← 百分比位置
    labeldistance=1.15,         # ← 標籤位置（餅圖外側）
    textprops={'fontsize': 11, 'color': 'white', 'weight': 'bold'},
    wedgeprops={'edgecolor': 'white', 'linewidth': 2}
)

# 百分比文字樣式
for autotext in autotexts:
    autotext.set_color('white')
    autotext.set_fontsize(12)
    autotext.set_weight('bold')

# 標籤文字樣式（深色，放在外側）
for text in texts:
    text.set_color('#1A1A1A')
    text.set_fontsize(11)
    text.set_weight('bold')
# ax.set_title('EPC Subsidiary Revenue Structure', fontsize=14, weight='bold', pad=20)
ax.set_title('EPC 子公司收入結構', fontsize=14, weight='bold', pad=20)
plt.tight_layout()
plt.savefig(f'{output_dir}chart_epc.png', dpi=150, bbox_inches='tight', facecolor='white')
plt.close()
print("[OK] chart_epc.png")

# ==================== 圖表 3：政策契合度雷達圖 ====================
fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(polar=True))
# categories = ['Computing-Energy', 'VPP', 'Green Power', 'Spot Market', 'REITs', 'Carbon', 'Zero-Carbon', 'V2G']
categories = ['算電協同', '虛擬電廠', '綠電直連', '電力現貨', '新能源REITs', '碳電協同', '零碳園區', 'V2G']
values = [95, 92, 90, 88, 85, 90, 87, 82]
values += values[:1]
angles = np.linspace(0, 2 * np.pi, len(categories), endpoint=False).tolist()
angles += angles[:1]

ax.plot(angles, values, 'o-', linewidth=2, color='#2B6CB4')
ax.fill(angles, values, alpha=0.25, color='#2B6CB4')
ax.set_xticks(angles[:-1])
ax.set_xticklabels(categories, fontsize=9)
ax.set_ylim(0, 100)
ax.set_yticks([60, 80, 100])
ax.set_yticklabels(['60', '80', '100'], fontsize=8, color='#999999')
# ax.set_title('Policy Alignment Score', fontsize=14, weight='bold', pad=20)
ax.set_title('政策契合度評分', fontsize=14, weight='bold', pad=20)
plt.tight_layout()
plt.savefig(f'{output_dir}chart_radar.png', dpi=150, bbox_inches='tight', facecolor='white')
plt.close()
print("[OK] chart_radar.png")

# ==================== 圖表 4：市場規模柱狀圖 ====================
fig, ax = plt.subplots(figsize=(8, 4.5))
# categories = ['VPP Market\n(Billion CNY)', 'Green Power\n(10k kW)', 'CCER Vol.\n(10k tons)', 'Carbon Price\n(CNY/ton)']
categories = ['VPP市場規模\n(億元)', '綠電直連裝機\n(萬kW)', 'CCER成交量\n(萬噸)', '碳價區間\n(元/噸)']
values = [1120, 340.5, 1580, 80]
colors_bar = ['#2B6CB4', '#2E8B57', '#D4692A', '#6B4C9A']
bars = ax.bar(categories, values, color=colors_bar, width=0.6, edgecolor='white', linewidth=1.5)
for bar, val in zip(bars, values):
    height = bar.get_height()
    ax.text(bar.get_x() + bar.get_width() / 2., height + max(values) * 0.02,
            f'{val}', ha='center', va='bottom', fontsize=11, weight='bold', color='#1A1A1A')
# ax.set_title('2026 China Energy Market Key Indicators', fontsize=14, weight='bold', pad=15)
ax.set_title('2026年中國能源市場關鍵指標', fontsize=14, weight='bold', pad=15)
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.spines['left'].set_color('#DDDDDD')
ax.spines['bottom'].set_color('#DDDDDD')
ax.tick_params(colors='#666666')
ax.set_ylabel('Value', fontsize=11, color='#666666')
plt.tight_layout()
plt.savefig(f'{output_dir}chart_market.png', dpi=150, bbox_inches='tight', facecolor='white')
plt.close()
print("[OK] chart_market.png")

# ==================== 圖表 5：商業閉環流程圖 ====================
fig, ax = plt.subplots(figsize=(12, 3))
ax.set_xlim(0, 12)
ax.set_ylim(0, 3)
ax.axis('off')

# boxes = [
#     (0.3, 'Trading Co.', 'Customer & Market Demand', '#FFF3E0', '#D4692A'),
#     (3.2, 'Investment/Dev Co.', 'Asset Control, VPP & AI', '#E8F5E9', '#2E8B57'),
#     (6.1, 'EPC Co.', 'Rapid Deployment', '#E3F2FD', '#2B6CB4'),
#     (9.0, 'Asset Pool', 'Scale Expansion & Returns', '#F3E5F5', '#6B4C9A'),
# ]
boxes = [(0.3, '貿易公司', '獲取客戶與市場需求', '#FFF3E0', '#D4692A'),
         (3.2, '投資開發公司', '控制資產、VPP與AI平台', '#E8F5E9', '#2E8B57'),
         (6.1, 'EPC公司', '快速部署與基建落地', '#E3F2FD', '#2B6CB4'),
         (9.0, '資產池', '規模擴張與收益增長', '#F3E5F5', '#6B4C9A')]

for x, t, d, fc, ec in boxes:
    rect = plt.Rectangle((x, 1), 2.2, 1.2, facecolor=fc, edgecolor=ec, linewidth=2, alpha=0.9)
    ax.add_patch(rect)
    ax.text(x + 1.1, 1.75, t, ha='center', va='center', fontsize=11, color=ec, weight='bold')
    ax.text(x + 1.1, 1.35, d, ha='center', va='center', fontsize=8, color='#666666')

arrow_style = dict(arrowstyle='->', color='#666666', lw=2)
for x in [2.6, 5.5, 8.4]:
    ax.annotate('', xy=(x + 0.5, 1.6), xytext=(x, 1.6), arrowprops=arrow_style)

# ax.set_title('Business Closed-Loop Flow', fontsize=14, weight='bold', pad=10)
ax.set_title('商業閉環流程', fontsize=14, weight='bold', pad=10)
plt.tight_layout()
plt.savefig(f'{output_dir}chart_loop.png', dpi=150, bbox_inches='tight', facecolor='white')
plt.close()
print("[OK] chart_loop.png")

# ==================== 圖表 6：風險評估矩陣 ====================
fig, ax = plt.subplots(figsize=(7, 5))
# risks = ['Price Volatility', 'Policy Diff.', 'CCER Method', 'Tech Integration', 'Capital Pressure']
risks = ['電力現貨價格波動', '區域政策差異', 'CCER方法學匹配', '技術整合', '資金沉澱']
impact = [4, 3, 3, 3, 4]
probability = [5, 4, 3, 2, 4]   # [4, 3, 3, 3, 4]
colors_risk = ['#C41E3A', '#D4692A', '#D4692A', '#D4692A', '#C41E3A']

ax.scatter(probability, impact, s=[300] * 5, c=colors_risk, alpha=0.7,
           edgecolors='white', linewidth=2, zorder=5)
for i, risk in enumerate(risks):
    ax.annotate(risk, (probability[i], impact[i]), textcoords="offset points",
                xytext=(0, 15), ha='center', fontsize=9, weight='bold')

ax.set_xlim(0.5, 5.5)
ax.set_ylim(0.5, 5.5)
# ax.set_xlabel('Probability', fontsize=11, color='#666666')
# ax.set_ylabel('Impact', fontsize=11, color='#666666')
# ax.set_title('Risk Assessment Matrix', fontsize=14, weight='bold', pad=15)
ax.set_xlabel('發生概率', fontsize=11, color='#666666')
ax.set_ylabel('影響程度', fontsize=11, color='#666666')
ax.set_title('風險評估矩陣', fontsize=14, weight='bold', pad=15)
ax.axhline(y=3, color='#DDDDDD', linestyle='--', linewidth=1)
ax.axvline(x=3, color='#DDDDDD', linestyle='--', linewidth=1)
ax.text(4.5, 4.8, 'High Attention', fontsize=10, color='#C41E3A', weight='bold', ha='center')
ax.text(1.5, 1.2, 'Low Attention', fontsize=10, color='#2E8B57', weight='bold', ha='center')
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.set_xticks([1, 2, 3, 4, 5])
ax.set_yticks([1, 2, 3, 4, 5])
ax.tick_params(colors='#666666')
plt.tight_layout()
plt.savefig(f'{output_dir}chart_risk.png', dpi=150, bbox_inches='tight', facecolor='white')
plt.close()
print("[OK] chart_risk.png")

# ==================== 開始構建 PPT ====================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ---------- 輔助函數 ----------
def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(6.5), Inches(2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    footer = slide.shapes.add_textbox(Inches(1), Inches(6.7), Inches(5), Inches(0.3))
    fp = footer.text_frame.paragraphs[0]
    fp.text = "Qianli Fengguang Plan - Organization Blueprint"
    fp.font.size = Pt(11)
    fp.font.color.rgb = LIGHT_GRAY
    return slide


def add_section_slide(prs, section_num, section_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT_BLUE
    bg.line.fill.background()

    num_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(2), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{section_num}" if section_num < 10 else str(section_num)
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1.2))
    tf2 = title_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = section_title
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    return slide


def add_content_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(10), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY
        p2.space_before = Pt(4)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.15), Inches(1.5), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    return slide


# ==================== SLIDE 1: 封面 ====================
add_title_slide(prs,
    "千里風光計畫",
    "行政區級 AI 光儲碳電網營運商組織藍圖\n與政策市場契合度研究報告")

# ==================== SLIDE 2: 目錄 ====================
slide = add_content_slide(prs, "目錄", "CONTENTS")
toc_items = [
    ("01", "千里風光計畫組織藍圖", "母公司定位 · 三大子公司 · 商業閉環 · 團隊建議"),
    ("02", "政策與市場契合度分析", "算電協同 · VPP與電力現貨 · 儲能市場化 · 碳交易 · 風險評估"),
]
y_pos = 1.6
for num, title, desc in toc_items:
    num_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(1), Inches(0.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    title_box = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(8), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK

    desc_box = slide.shapes.add_textbox(Inches(2), Inches(y_pos + 0.4), Inches(8), Inches(0.3))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY

    y_pos += 1.5

# ==================== SLIDE 3: 第一部分章節頁 ====================
add_section_slide(prs, 1, "千里風光計畫組織藍圖")

# ==================== SLIDE 4: 摘要與核心願景 ====================
slide = add_content_slide(prs, "摘要與核心願景", "從能源建設商轉型為城市級能源大腦")

# 左側卡片
box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(6), Inches(2.5))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
box1.line.color.rgb = ACCENT_BLUE
box1.line.width = Pt(1.5)

tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(5.4), Inches(2))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "計畫定位"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

p2 = tf.add_paragraph()
p2.text = "建立以 AI 為核心驅動的行政區級光伏、儲能與碳電網營運商"
p2.font.size = Pt(13)
p2.font.color.rgb = DARK
p2.space_before = Pt(8)

p3 = tf.add_paragraph()
p3.text = "核心願景：能源資產投資平台、區域電網營運商、能源金融集團"
p3.font.size = Pt(13)
p3.font.color.rgb = DARK
p3.space_before = Pt(6)

# 右側卡片
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.5), Inches(5.7), Inches(2.5))
box2.fill.solid()
box2.fill.fore_color.rgb = WHITE
box2.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

caps = [
    ("資產控制", "透過 EPC 快速部署獲取區域能源資產控制權"),
    ("資本調度", "設計 REITs 2.0 退出路徑，實現資本循環"),
    ("數據中樞", "AI 能源平台匯聚全域能源數據，驅動智慧決策"),
    ("交易權與政府接口", "掌握電力現貨、碳交易、綠電直連等市場入口"),
]
y = 1.7
for cap_title, cap_desc in caps:
    tb = slide.shapes.add_textbox(Inches(7.3), Inches(y), Inches(5.1), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"● {cap_title}：{cap_desc}"
    p.font.size = Pt(11)
    p.font.color.rgb = DARK
    y += 0.35

# 底部三大前瞻佈局
box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.3), Inches(12.1), Inches(1.8))
box3.fill.solid()
box3.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
box3.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

tb = slide.shapes.add_textbox(Inches(0.9), Inches(4.5), Inches(11.5), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "三大前瞻佈局"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK

layouts = [
    ("算電協同", "數據中心等高耗能負載與新能源發電深度耦合", ACCENT_BLUE),
    ("多用戶綠電直連", "為園區內多個用戶提供專線綠電供應", ACCENT_GREEN),
    ("V2G 行動儲能池", "將電動車輛視為可調度的移動儲能節點", ACCENT_ORANGE),
]
x = 0.9
for lt, ld, lc in layouts:
    tb = slide.shapes.add_textbox(Inches(x), Inches(4.9), Inches(3.8), Inches(1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = lt
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = lc

    p2 = tf.add_paragraph()
    p2.text = ld
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(4)
    x += 4.1

# ==================== SLIDE 5: 集團母公司定位 ====================
slide = add_content_slide(prs, "集團母公司定位", "區域級智慧電網資產控股與營運平台")

funcs = [
    ("資產控制", "戰略規劃與資源整合\n掌握關鍵能源資產", ACCENT_BLUE),
    ("資本調度", "資本靈活調度\nREITs 退出路徑設計", ACCENT_GREEN),
    ("數據中樞", "能源數據匯聚\nAI 智慧決策中樞", ACCENT_ORANGE),
    ("交易權與政府接口", "市場交易權掌握\n政府關係與政策接口", ACCENT_PURPLE),
]
x = 0.6
for title, desc, color in funcs:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(2.8), Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(1.8), Inches(2.5), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color

    tb2 = slide.shapes.add_textbox(Inches(x + 0.15), Inches(2.3), Inches(2.5), Inches(1.2))
    tf = tb2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY

    x += 3.1

# 未來發展方向
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.5))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
box.line.color.rgb = ACCENT_BLUE
box.line.width = Pt(1)

tb = slide.shapes.add_textbox(Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "未來發展方向"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK

futures = [
    ("算電協同", "將數據中心等高耗能負載與新能源發電深度耦合，實現能源效率最大化", ACCENT_BLUE),
    ("多用戶綠電直連", "為園區內多個用戶提供專線綠電供應，提升綠電消納能力", ACCENT_GREEN),
    ("V2G 行動儲能池", "將行政區內電動車輛視為可調度的移動儲能節點，為電網提供靈活性資源", ACCENT_ORANGE),
]
y = 4.9
for ft, fd, fc in futures:
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(y), Inches(11.5), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"● "
    p.font.size = Pt(12)
    p.font.color.rgb = fc
    p.font.bold = True

    run = p.add_run()
    run.text = f"{ft}："
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = fc

    run2 = p.add_run()
    run2.text = fd
    run2.font.size = Pt(12)
    run2.font.color.rgb = GRAY
    y += 0.4

# ==================== SLIDE 6: 組織架構全景 ====================
slide = add_content_slide(prs, "母子公司組織架構全景", "母公司統籌 + 三大子公司專業協作")

slide.shapes.add_picture(f'{output_dir}chart_org.png', Inches(1.5), Inches(1.5), width=Inches(10.3))

box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.8))
box.fill.solid()
box.fill.fore_color.rgb = WHITE
box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

collabs = [
    ("EPC 工程總包", "用工程能力拿下區域資產控制權", ACCENT_BLUE),
    ("投資/開發/方案", "持有資產、控制 VPP、管理電網收益", ACCENT_GREEN),
    ("貿易公司", "市場擴張、獲取客戶、輸出方案", ACCENT_ORANGE),
]
x = 0.9
for ct, cd, cc in collabs:
    tb = slide.shapes.add_textbox(Inches(x), Inches(5.4), Inches(3.8), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ct
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = cc

    p2 = tf.add_paragraph()
    p2.text = cd
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(4)
    x += 4.1

# ==================== SLIDE 7: EPC 子公司 ====================
slide = add_content_slide(prs, "子公司 1：EPC 工程總包公司", "行政區能源基建快速部署平台")

box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(6), Inches(2.8))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
box1.line.color.rgb = ACCENT_BLUE

tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(5.4), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "定位與核心業務"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

tb2 = slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(5.4), Inches(1.8))
tf = tb2.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "提供標準化、模塊化、可快速複製的工程能力，而非僅僅賺取工程毛利"
p.font.size = Pt(12)
p.font.color.rgb = DARK

services = ["光伏與儲能設計施工", "微電網建設", "Edge AI 設備安裝", "系統運維與舊系統改造"]
for s in services:
    p = tf.add_paragraph()
    p.text = f"  • {s}"
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.space_before = Pt(4)

box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.5), Inches(5.7), Inches(2.8))
box2.fill.solid()
box2.fill.fore_color.rgb = WHITE
box2.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

tb = slide.shapes.add_textbox(Inches(7.3), Inches(1.7), Inches(5.1), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "組織架構與「戰區制」管理"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = DARK

deps = ["工程中心", "調試中心", "O&M 中心", "AI 接入部", "採購部"]
y = 2.1
for d in deps:
    tb = slide.shapes.add_textbox(Inches(7.3), Inches(y), Inches(2.3), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"• {d}"
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    y += 0.3

box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(3.3), Inches(5.1), Inches(0.8))
box3.fill.solid()
box3.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
box3.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

tb = slide.shapes.add_textbox(Inches(7.5), Inches(3.45), Inches(4.7), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "管理模式：SOP 化 · 模塊化 · 數字工地 · 區域責任制"
p.font.size = Pt(11)
p.font.color.rgb = DARK

# 收入圖表
slide.shapes.add_picture(f'{output_dir}chart_epc.png', Inches(0.6), Inches(4.6), width=Inches(3)) # width=Inches(4.5)

box4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(4.6), Inches(7.2), Inches(2.5)) # Inches(2.2)
box4.fill.solid()
box4.fill.fore_color.rgb = WHITE
box4.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

tb = slide.shapes.add_textbox(Inches(5.8), Inches(4.8), Inches(6.6), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "技術專業能力"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = DARK

techs = ["高低壓配電、PCS/BMS", "EMS 能源管理系統", "Modbus / IEC104 通訊協議",
         "SCADA 監控系統", "Edge AI 本地控制", "數字孿生技術"]
y = 5.2
for t in techs:
    tb = slide.shapes.add_textbox(Inches(5.8), Inches(y), Inches(3), Inches(0.25))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"• {t}"
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    y += 0.28

# ==================== SLIDE 8: 投資開發公司 ====================
slide = add_content_slide(prs, "子公司 2：投資 / 開發 / 方案公司", "集團真正核心 — 行政區能源資產運營商")

box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
box1.line.color.rgb = ACCENT_GREEN
box1.line.width = Pt(2)

tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.6))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "負責持有資產、控制 VPP、管理電網收益並控制現金流。核心目標是管理整個行政區的能源資產池，而非單一案場，以實現資產收益的最大化。"
p.font.size = Pt(13)
p.font.color.rgb = DARK

box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.7), Inches(6), Inches(3.8))
box2.fill.solid()
box2.fill.fore_color.rgb = WHITE
box2.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.9), Inches(5.4), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "六大核心部門"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = DARK

deps = [
    ("投資中心", "項目投資決策與資產配置"),
    ("VPP 中心", "虛擬電廠聚合與調度"),
    ("AI 中心", "算法研發與平台運營"),
    ("能源交易部", "電力現貨與中長期交易"),
    ("碳資產部", "CCER 開發與碳交易"),
    ("財務資產部", "資產證券化與 REITs"),
]
y = 3.3
for dtitle, ddesc in deps:
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(y), Inches(5.4), Inches(0.45))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"● {dtitle}："
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    run = p.add_run()
    run.text = ddesc
    run.font.size = Pt(11)
    run.font.color.rgb = GRAY
    y += 0.35

box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2.7), Inches(5.7), Inches(2))
box3.fill.solid()
box3.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
box3.line.color.rgb = ACCENT_BLUE

tb = slide.shapes.add_textbox(Inches(7.3), Inches(2.9), Inches(5.1), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "四大現金流來源"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = DARK

sources = [
    ("售電", "穩定現金流", ACCENT_BLUE),
    ("VPP 調度", "高毛利", ACCENT_GREEN),
    ("碳交易", "金融收益", ACCENT_ORANGE),
    ("AI SaaS", "複利增長", ACCENT_PURPLE),
]
y = 3.3
for sname, sdesc, sc in sources:
    tb = slide.shapes.add_textbox(Inches(7.3), Inches(y), Inches(5.1), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"● {sname}（{sdesc}）"
    p.font.size = Pt(11)
    p.font.color.rgb = sc
    p.font.bold = True
    y += 0.32

box4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(4.9), Inches(5.7), Inches(1.6))
box4.fill.solid()
box4.fill.fore_color.rgb = WHITE
box4.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

tb = slide.shapes.add_textbox(Inches(7.3), Inches(5.1), Inches(5.1), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "關鍵績效指標 (KPI)"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = DARK

kpis = ["聚合 MW / MWh 規模", "調度收益總額", "每 MW 資產收益"]
y = 5.45
for k in kpis:
    tb = slide.shapes.add_textbox(Inches(7.3), Inches(y), Inches(5.1), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"• {k}"
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    y += 0.28

# ==================== SLIDE 9: 貿易公司 ====================
slide = add_content_slide(prs, "子公司 3：貿易公司", "能源資產商品化輸出平台 — 能源市場擴張引擎")

box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.8))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
box1.line.color.rgb = ACCENT_ORANGE
box1.line.width = Pt(2)

tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.65), Inches(11.5), Inches(0.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "負責市場擴張、區域滲透及潛在的海外輸出。將能源資產與服務商品化，輸出至終端市場。"
p.font.size = Pt(13)
p.font.color.rgb = DARK

box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.5), Inches(6), Inches(2.2))
box2.fill.solid()
box2.fill.fore_color.rgb = WHITE
box2.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.7), Inches(5.4), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "核心業務"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = DARK

bizs = ["售電方案", "光儲方案", "AI 能源平台 SaaS", "VPP 加盟合作"]
y = 3.1
for b in bizs:
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(y), Inches(5.4), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"• {b}"
    p.font.size = Pt(12)
    p.font.color.rgb = ACCENT_ORANGE
    p.font.bold = True
    y += 0.35

box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2.5), Inches(5.7), Inches(2.2))
box3.fill.solid()
box3.fill.fore_color.rgb = WHITE
box3.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

tb = slide.shapes.add_textbox(Inches(7.3), Inches(2.7), Inches(5.1), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "「行業制」管理"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = DARK

industries = ["工業", "物流", "數據中心", "港口"]
y = 3.1
for ind in industries:
    tb = slide.shapes.add_textbox(Inches(7.3), Inches(y), Inches(2.3), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"• {ind}"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    y += 0.35

box4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.9), Inches(12.1), Inches(1.8))
box4.fill.solid()
box4.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
box4.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

tb = slide.shapes.add_textbox(Inches(0.9), Inches(5.1), Inches(11.5), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "市場拓展建議"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = DARK

suggestions = [
    ("供應鏈碳足跡諮詢", "對接歐盟 CBAM 等國際標準，幫助出口企業應對碳關稅", ACCENT_BLUE),
    ("「能源+」增值服務", "能效診斷、設備預警 SaaS 等增值服務", ACCENT_GREEN),
    ("區域能源合夥人計畫", "招募地方資源方，快速滲透縣域市場", ACCENT_ORANGE),
]
x = 0.9
for st, sd, sc in suggestions:
    tb = slide.shapes.add_textbox(Inches(x), Inches(5.5), Inches(3.8), Inches(1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = st
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = sc

    p2 = tf.add_paragraph()
    p2.text = sd
    p2.font.size = Pt(10)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(3)
    x += 4.1

# ==================== SLIDE 10: 商業閉環 ====================
slide = add_content_slide(prs, "商業閉環與最終形態", "貿易獲客 → 投資控資產 → EPC 落地 → 資產池擴張")

slide.shapes.add_picture(f'{output_dir}chart_loop.png', Inches(0.6), Inches(1.5), width=Inches(12.1))

box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.6), Inches(12.3), Inches(2.4))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

tb = slide.shapes.add_textbox(Inches(0.9), Inches(4.8), Inches(11.5), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "最終形態演變"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK

forms = [
    ("區域電網營運商", "管理行政區級能源資產池，參與電力市場交易", ACCENT_BLUE),
    ("能源金融平台", "REITs、碳金融、綠色金融產品創新", ACCENT_GREEN),
    ("AI 能源 OS", "城市級能源大腦，數據驅動的智慧決策中樞", ACCENT_PURPLE),
]
x = 0.9
for ft, fd, fc in forms:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.3), Inches(3.8), Inches(1.4))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = fc
    card.line.width = Pt(2)

    tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(5.5), Inches(3.5), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = ft
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = fc

    tb2 = slide.shapes.add_textbox(Inches(x + 0.15), Inches(5.95), Inches(3.5), Inches(0.6))
    tf = tb2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = fd
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY

    x += 4       # 4.1

# ==================== SLIDE 11: 估值與團隊 ====================
slide = add_content_slide(prs, "未來估值與團隊組成建議", "從 EPC 估值轉向資產平台估值與 AI+VPP+碳金融估值")

box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(6), Inches(3))
box1.fill.solid()
box1.fill.fore_color.rgb = WHITE
box1.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(5.4), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "估值邏輯轉變"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = DARK

vals = [
    ("傳統 EPC 估值", "基於施工毛利與工程規模", GRAY),
    ("↓", "", GRAY),
    ("資產平台估值", "基於持續性電網現金流、服務分潤", ACCENT_BLUE),
    ("↓", "", GRAY),
    ("AI + VPP + 碳金融估值", "基於數據資產溢價與市場化收益", ACCENT_GREEN),
]
y = 2.1
for vt, vd, vc in vals:
    if vt == "↓":
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(y), Inches(5.4), Inches(0.2))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = vt
        p.font.size = Pt(14)
        p.font.color.rgb = vc
        p.alignment = PP_ALIGN.CENTER
    else:
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(y), Inches(5.4), Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = vt
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = vc
        if vd:
            p2 = tf.add_paragraph()
            p2.text = vd
            p2.font.size = Pt(11)
            p2.font.color.rgb = GRAY
    y += 0.35

box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.5), Inches(5.7), Inches(3))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
box2.line.color.rgb = ACCENT_BLUE

tb = slide.shapes.add_textbox(Inches(7.3), Inches(1.7), Inches(5.1), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "跨界人才結構"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = DARK

tb2 = slide.shapes.add_textbox(Inches(7.3), Inches(2.1), Inches(5.1), Inches(0.3))
tf = tb2.text_frame
p = tf.paragraphs[0]
p.text = "電力系統 + AI 算法 + 金融工程"
p.font.size = Pt(11)
p.font.color.rgb = GRAY

roles = [
    ("能源交易員", "電力現貨與碳市場交易", ACCENT_BLUE),
    ("AI 調度算法工程師", "VPP 優化與預測算法", ACCENT_GREEN),
    ("碳資產架構師", "CCER 開發與碳金融設計", ACCENT_ORANGE),
]
y = 2.5
for rt, rd, rc in roles:
    tb = slide.shapes.add_textbox(Inches(7.3), Inches(y), Inches(5.1), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"● {rt}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = rc

    p2 = tf.add_paragraph()
    p2.text = f"  {rd}"
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(2)
    y += 0.45

box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.8))
box3.fill.solid()
box3.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
box3.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

tb = slide.shapes.add_textbox(Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "合夥人激勵機制"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = DARK

incentives = [
    ("項目跟投", "核心骨干參與項目投資", ACCENT_BLUE),
    ("期權激勵", "長期股權綁定", ACCENT_GREEN),
    ("資產池收益分潤", "與行政區長期收益掛鉤", ACCENT_ORANGE),
]
x = 0.9
for it, idesc, ic in incentives:
    tb = slide.shapes.add_textbox(Inches(x), Inches(5.4), Inches(3.8), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = it
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ic

    p2 = tf.add_paragraph()
    p2.text = idesc
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(3)
    x += 4.1

# ==================== SLIDE 12: 結論 ====================
slide = add_content_slide(prs, "結論", "搶佔中國新能源轉型浪潮先機")

box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.2))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
box1.line.color.rgb = ACCENT_GREEN
box1.line.width = Pt(2)

tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "「千里風光計畫」提供了一個全面且具前瞻性的戰略框架，透過清晰的母子公司定位、精細化的營運管理、創新的財務模式與強大的技術支撐，有望在行政區級市場建立起領先的 AI 光儲碳電網營運商。"
p.font.size = Pt(14)
p.font.color.rgb = DARK

box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.9), Inches(6), Inches(3.5))
box2.fill.solid()
box2.fill.fore_color.rgb = WHITE
box2.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

tb = slide.shapes.add_textbox(Inches(0.9), Inches(3.1), Inches(5.4), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "核心優勢"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK

advantages = [
    "清晰的組織架構：母公司統籌 + 三大子公司專業分工",
    "完整的商業閉環：從獲客到資產退出的全鏈條覆蓋",
    "AI 驅動的技術護城河：Edge AI、數字孿生、智能調度",
    "金融創新能力：REITs、碳交易、電力現貨聯動",
]
y = 3.5
for adv in advantages:
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(y), Inches(5.4), Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"✓ {adv}"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK
    y += 0.4

box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2.9), Inches(5.7), Inches(3.5))
box3.fill.solid()
box3.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
box3.line.color.rgb = ACCENT_BLUE

tb = slide.shapes.add_textbox(Inches(7.3), Inches(3.1), Inches(5.1), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "行動建議"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK

actions = [
    "儘快啟動核心團隊組建，鎖定電力+AI+金融複合人才",
    "選定首個行政區試點，驗證「戰區制」與商業閉環",
    "搭建 AI 能源平台雛形，積累數據資產與算法能力",
    "設計合夥人激勵機制，綁定核心骨干長期利益",
]
y = 3.5
for act in actions:
    tb = slide.shapes.add_textbox(Inches(7.3), Inches(y), Inches(5.1), Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"→ {act}"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK
    y += 0.4

box4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.6))
box4.fill.solid()
box4.fill.fore_color.rgb = DARK
box4.line.fill.background()

tb = slide.shapes.add_textbox(Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "最終願景：城市級能源大腦  |  能源建設商 → 能源資產投資平台 → 區域電網營運商 → 能源金融集團"
p.font.size = Pt(12)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# ==================== 保存 PPT ====================
output_path = f'{output_dir}千里風光計畫_完整簡報.pptx'
prs.save(output_path)
print(f"\n{'='*50}")
print(f"[SUCCESS] PPT 已保存至: {output_path}")
print(f"{'='*50}")
print(f"共生成 {len(prs.slides)} 頁幻燈片")
print(f"共生成 6 張圖表於: {output_dir}")