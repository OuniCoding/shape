from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import matplotlib
import numpy as np
import os



matplotlib.use('Agg')
# plt.rcParams['font.sans-serif'] = ['SimHei', 'DejaVu Sans']
# 微軟正黑體或微軟雅黑
plt.rcParams['font.sans-serif'] = ['Microsoft YaHei', 'SimHei']
plt.rcParams['axes.unicode_minus'] = False


output_dir = './ppt_output/'
os.makedirs(output_dir, exist_ok=True)

# ========== 顏色定義 ==========
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_BLUE = RGBColor(0x2B, 0x6C, 0xB4)
ACCENT_GREEN = RGBColor(0x2E, 0x8B, 0x57)
ACCENT_ORANGE = RGBColor(0xD4, 0x69, 0x2A)
ACCENT_PURPLE = RGBColor(0x6B, 0x4C, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ========== 生成圖表 ==========
# 圖表1: 組織架構圖
fig, ax = plt.subplots(figsize=(10, 5))
ax.set_xlim(0, 10); ax.set_ylim(0, 5); ax.axis('off')
rect1 = plt.Rectangle((3.5, 3.8), 3, 0.8, facecolor='#2B6CB4', edgecolor='#1A1A1A', linewidth=2)
ax.add_patch(rect1)
ax.text(5, 4.2, '集團母公司\n區域級智慧電網資產控股與營運平台', ha='center', va='center', fontsize=11, color='white', weight='bold')
ax.plot([5, 5], [3.8, 3.2], color='#666666', linewidth=2)
ax.plot([1.5, 8.5], [3.2, 3.2], color='#666666', linewidth=2)
for x_pos, title, color, desc in [(1.5, 'EPC 工程總包公司', '#2B6CB4', '區域工程交付平台'),
                                   (5, '投資/開發/方案公司', '#2E8B57', '集團真正核心'),
                                   (8.5, '貿易公司', '#D4692A', '能源市場擴張引擎')]:
    ax.plot([x_pos, x_pos], [3.2, 2.8], color='#666666', linewidth=2)
    rect = plt.Rectangle((x_pos-1.3, 1.8), 2.6, 0.9, facecolor={'#2B6CB4':'#E3F2FD','#2E8B57':'#E8F5E9','#D4692A':'#FFF3E0'}[color],
                          edgecolor=color, linewidth=2)
    ax.add_patch(rect)
    ax.text(x_pos, 2.35, f'子公司\n{title}', ha='center', va='center', fontsize=10, color=color, weight='bold')
plt.tight_layout()
plt.savefig(f'{output_dir}chart_org.png', dpi=150, bbox_inches='tight', facecolor='white')
plt.close()

# 圖表2: EPC收入結構
fig, ax = plt.subplots(figsize=(6, 4))
labels = ['EPC 施工\n40%', '運維\n30%', '改造升級\n20%', 'AI 部署\n10%']
sizes = [40, 30, 20, 10]
colors = ['#2B6CB4', '#2E8B57', '#D4692A', '#6B4C9A']
ax.pie(sizes, labels=labels, colors=colors, startangle=90, textprops={'fontsize': 10, 'color': 'white', 'weight': 'bold'})
ax.set_title('EPC 子公司收入結構', fontsize=14, weight='bold', pad=20)
plt.tight_layout()
plt.savefig(f'{output_dir}chart_epc.png', dpi=150, bbox_inches='tight', facecolor='white')
plt.close()

# 圖表3: 政策雷達圖
fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(polar=True))
categories = ['算電協同', '虛擬電廠', '綠電直連', '電力現貨', '新能源REITs', '碳電協同', '零碳園區', 'V2G']
values = [95, 92, 90, 88, 85, 90, 87, 82]
values += values[:1]
angles = np.linspace(0, 2*np.pi, len(categories), endpoint=False).tolist()
angles += angles[:1]
ax.plot(angles, values, 'o-', linewidth=2, color='#2B6CB4')
ax.fill(angles, values, alpha=0.25, color='#2B6CB4')
ax.set_xticks(angles[:-1]); ax.set_xticklabels(categories, fontsize=10)
ax.set_ylim(0, 100); ax.set_title('政策契合度評分', fontsize=14, weight='bold', pad=20)
plt.tight_layout()
plt.savefig(f'{output_dir}chart_radar.png', dpi=150, bbox_inches='tight', facecolor='white')
plt.close()

# 圖表4: 市場規模
fig, ax = plt.subplots(figsize=(8, 4.5))
categories = ['VPP市場規模\n(億元)', '綠電直連裝機\n(萬kW)', 'CCER成交量\n(萬噸)', '碳價區間\n(元/噸)']
values = [1120, 340.5, 1580, 80]
colors_bar = ['#2B6CB4', '#2E8B57', '#D4692A', '#6B4C9A']
bars = ax.bar(categories, values, color=colors_bar, width=0.6, edgecolor='white', linewidth=1.5)
for bar, val in zip(bars, values):
    ax.text(bar.get_x()+bar.get_width()/2., bar.get_height()+max(values)*0.02, str(val),
            ha='center', va='bottom', fontsize=11, weight='bold', color='#1A1A1A')
ax.set_title('2026年中國能源市場關鍵指標', fontsize=14, weight='bold', pad=15)
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
plt.tight_layout()
plt.savefig(f'{output_dir}chart_market.png', dpi=150, bbox_inches='tight', facecolor='white')
plt.close()

# 圖表5: 商業閉環
fig, ax = plt.subplots(figsize=(12, 3))
ax.set_xlim(0, 12); ax.set_ylim(0, 3); ax.axis('off')
boxes = [(0.3, '貿易公司', '獲取客戶與市場需求', '#FFF3E0', '#D4692A'),
         (3.2, '投資開發公司', '控制資產、VPP與AI平台', '#E8F5E9', '#2E8B57'),
         (6.1, 'EPC公司', '快速部署與基建落地', '#E3F2FD', '#2B6CB4'),
         (9.0, '資產池', '規模擴張與收益增長', '#F3E5F5', '#6B4C9A')]
for x, t, d, fc, ec in boxes:
    rect = plt.Rectangle((x, 1), 2.2, 1.2, facecolor=fc, edgecolor=ec, linewidth=2)
    ax.add_patch(rect)
    ax.text(x+1.1, 1.75, t, ha='center', va='center', fontsize=11, color=ec, weight='bold')
    ax.text(x+1.1, 1.35, d, ha='center', va='center', fontsize=8, color='#666666')
for x in [2.6, 5.5, 8.4]:
    ax.annotate('', xy=(x+0.5, 1.6), xytext=(x, 1.6), arrowprops=dict(arrowstyle='->', color='#666666', lw=2))
ax.set_title('商業閉環流程', fontsize=14, weight='bold', pad=10)
plt.tight_layout()
plt.savefig(f'{output_dir}chart_loop.png', dpi=150, bbox_inches='tight', facecolor='white')
plt.close()

# 圖表6: 風險矩陣
fig, ax = plt.subplots(figsize=(7, 5))
risks = ['電力現貨價格波動', '區域政策差異', 'CCER方法學匹配', '技術整合', '資金沉澱']
impact = [4, 3, 3, 3, 4]
probability = [5, 4, 3, 2, 4]# [4, 3, 3, 3, 4]
colors_risk = ['#C41E3A', '#D4692A', '#D4692A', '#D4692A', '#C41E3A']
ax.scatter(probability, impact, s=[300]*5, c=colors_risk, alpha=0.7, edgecolors='white', linewidth=2, zorder=5)
for i, risk in enumerate(risks):
    ax.annotate(risk, (probability[i], impact[i]), textcoords="offset points", xytext=(0, 15),
                ha='center', fontsize=9, weight='bold')
ax.set_xlim(0.5, 5.5); ax.set_ylim(0.5, 5.5)
ax.set_xlabel('發生概率', fontsize=11, color='#666666')
ax.set_ylabel('影響程度', fontsize=11, color='#666666')
ax.set_title('風險評估矩陣', fontsize=14, weight='bold', pad=15)
ax.axhline(y=3, color='#DDDDDD', linestyle='--', linewidth=1)
ax.axvline(x=3, color='#DDDDDD', linestyle='--', linewidth=1)
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
plt.tight_layout()
plt.savefig(f'{output_dir}chart_risk.png', dpi=150, bbox_inches='tight', facecolor='white')
plt.close()

print("所有圖表生成完成！")

# ========== 創建 PPT ==========
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    p = tb.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    p2 = tb2.text_frame.paragraphs[0]; p2.text = subtitle; p2.font.size = Pt(18); p2.font.color.rgb = LIGHT_GRAY
    return slide

def add_section_slide(prs, num, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = ACCENT_BLUE; bg.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(2), Inches(1))
    p = tb.text_frame.paragraphs[0]; p.text = f"0{num}"; p.font.size = Pt(72); p.font.bold = True; p.font.color.rgb = RGBColor(0xCC,0xCC,0xCC)
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1.2))
    p2 = tb2.text_frame.paragraphs[0]; p2.text = title; p2.font.size = Pt(36); p2.font.bold = True; p2.font.color.rgb = WHITE
    return slide

def add_content_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(10), Inches(0.6))
    p = tb.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = DARK
    if subtitle:
        p2 = tb.text_frame.add_paragraph(); p2.text = subtitle; p2.font.size = Pt(12); p2.font.color.rgb = GRAY; p2.space_before = Pt(4)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.15), Inches(1.5), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT_BLUE; line.line.fill.background()
    return slide

# [此處插入前面已構建的12頁內容代碼...]
# 由於篇幅限制，建議將前面 ipython 中的幻燈片構建代碼複製至此處

# 保存檔案
prs.save(f'{output_dir}千里風光計畫_完整簡報.pptx')
print(f"PPT 已保存至: {output_dir}千里風光計畫_完整簡報.pptx")